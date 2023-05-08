import os
import time
import cProfile

import numpy as np

from pptx import Presentation
from pptx.util import Inches

from roshambo.core import GetSimilarityScores
from roshambo.analysis import (
    calc_roc_auc,
    plot_mult_roc,
    plot_mult_auc,
    plot_mult_enrichment,
)

main_dir = "/UserUCDD/ratwi/tests/dud-z"

for folder in os.listdir(main_dir):
    if folder in ["UROK_new_DUDE_1"]:
        try:
            for subfolder in ["ligands", "decoys"]:
                curr_dir = f"{main_dir}/{folder}/{subfolder}"
                os.chdir(curr_dir)
                print(curr_dir)
                st = time.time()
                sm = GetSimilarityScores(
                    ref_file="../xtal-lig.sdf",
                    dataset_files_pattern=f"{subfolder}.smi",
                    ignore_hs=True,
                    n_confs=50,
                    keep_mol=True,
                    random_seed=109838974,
                    opt_confs=False,
                    calc_energy=False,
                    energy_iters=300,
                    energy_cutoff=np.inf,
                    align_confs=True,
                    rms_cutoff=0.1,
                    num_threads=28,
                    method="ETKDGv2",
                    working_dir=None,
                    # smiles_kwargs={"delimiter": "\t"},
                )
                sm.run_paper()
                sm.convert_transformation_arrays()
                sm.transform_molecules()
                sm.calculate_scores(
                    volume_type="analytic",
                    n=2,
                    epsilon=0.5,
                    use_carbon_radii=True,
                    color=True,
                    max_conformers=1,
                    sort_by="ComboTanimoto",
                    write_to_file=True,
                )
                et = time.time()
                print(f"Full calculation took: {et - st}")
                os.chdir(main_dir)
        except Exception as e:
            print(e)

filtered_folders = [
    "RENI_new_DUDE_1",
    "ROCK1_new_DUDE_1",
    "SRC_new_DUDE_1",
    "THRB_new_DUDE_1",
    "TRY1_new_DUDE_1",
    "TRYB1_new_DUDE_1",
    "UROK_new_DUDE_1",
    "XIAP_new_DUDE_1",
]
for folder in os.listdir(main_dir):
    if folder in filtered_folders:
        try:
            curr_dir = f"{main_dir}/{folder}"
            os.chdir(curr_dir)
            print(curr_dir)
            calc_roc_auc(
                f"{curr_dir}/ligands/roshambo.csv",
                f"{curr_dir}/decoys/roshambo.csv",
                score="ComboTanimoto",
                n_bootstraps=1000,
                interpolation=True,
            )
            os.rename("analysis.csv", f"pypaper_{folder}_analysis.csv")
            os.rename("roc.csv", f"pypaper_{folder}_roc.csv")
            os.rename("auc_roc.jpg", f"pypaper_{folder}_auc_roc.jpg")
            os.chdir(main_dir)
        except Exception as e:
            print(e)

pypaper_analysis = [
    f"{main_dir}/{name}/pypaper_{name}_analysis.csv" for name in filtered_folders
]
pypaper_roc = [f"{main_dir}/{name}/pypaper_{name}_roc.csv" for name in filtered_folders]
rocs_analysis = [
    f"{main_dir}/{name}/rocs_{name}_analysis.csv" for name in filtered_folders
]
rocs_roc = [f"{main_dir}/{name}/rocs_{name}_roc.csv" for name in filtered_folders]

for name, roc_p, analysis_p, roc_r, analysis_r in zip(
    filtered_folders, pypaper_roc, pypaper_analysis, rocs_roc, rocs_analysis
):
    plot_mult_roc(
        rates_dict={"PYPAPER": roc_p, "ROCS": roc_r},
        analysis_dict={"PYPAPER": analysis_p, "ROCS": analysis_r},
        colors_dict={"PYPAPER": "#807FFF", "ROCS": "#7FC080"},
        title=f"{name} ROC",
        log=False,
        filename=f"{name}_roc.jpg",
    )
plot_mult_auc(
    auc_dict={"PYPAPER": pypaper_analysis, "ROCS": rocs_analysis},
    colors_dict={"PYPAPER": "#807FFF", "ROCS": "#7FC080"},
    group_labels=[i + 36 for i in range(len(filtered_folders))],
    # figsize=(12, 8),
)
plot_mult_enrichment(
    enrich_dict={"PYPAPER": pypaper_analysis, "ROCS": rocs_analysis},
    colors_dict={0: "#807FFF", 1: "#7FC080", 2: "gray", 3: "black"},
    hatch_patterns=[None, "+"],
    group_labels=[i + 36 for i in range(len(filtered_folders))],
    # figsize=(12, 8),
)

import pandas as pd

auc_diff = []
enrich_1 = []
enrich_2 = []
enrich_3 = []
enrich_4 = []
for folder in os.listdir(main_dir):
    try:
        print(folder)
        df_pypaper = pd.read_csv(
            f"{main_dir}/{folder}/pypaper_{folder}_analysis.csv", delimiter="\t"
        )
        df_rocs = pd.read_csv(
            f"{main_dir}/{folder}/rocs_{folder}_analysis.csv", delimiter="\t"
        )

        rocs_auc = df_rocs.loc[df_rocs["Run Name"] == "AUC", "Mean"].values[0]
        pypaper_auc = df_pypaper.loc[df_pypaper["Run Name"] == "AUC", "Mean"].values[0]
        auc_diff.append((rocs_auc - pypaper_auc) / rocs_auc)

        rocs_enrich_1 = df_rocs.loc[
            df_rocs["Run Name"] == "0.5% Enrichment", "Mean"
        ].values[0]
        pypaper_enrich_1 = df_pypaper.loc[
            df_pypaper["Run Name"] == "0.5% Enrichment", "Mean"
        ].values[0]
        enrich_1.append((rocs_enrich_1 - pypaper_enrich_1) / rocs_enrich_1)

        rocs_enrich_2 = df_rocs.loc[
            df_rocs["Run Name"] == "1.0% Enrichment", "Mean"
        ].values[0]
        pypaper_enrich_2 = df_pypaper.loc[
            df_pypaper["Run Name"] == "1.0% Enrichment", "Mean"
        ].values[0]
        enrich_2.append((rocs_enrich_2 - pypaper_enrich_2) / rocs_enrich_2)

        rocs_enrich_3 = df_rocs.loc[
            df_rocs["Run Name"] == "2.0% Enrichment", "Mean"
        ].values[0]
        pypaper_enrich_3 = df_pypaper.loc[
            df_pypaper["Run Name"] == "2.0% Enrichment", "Mean"
        ].values[0]
        enrich_3.append((rocs_enrich_3 - pypaper_enrich_3) / rocs_enrich_3)

        rocs_enrich_4 = df_rocs.loc[
            df_rocs["Run Name"] == "5.0% Enrichment", "Mean"
        ].values[0]
        pypaper_enrich_4 = df_pypaper.loc[
            df_pypaper["Run Name"] == "5.0% Enrichment", "Mean"
        ].values[0]
        enrich_4.append((rocs_enrich_4 - pypaper_enrich_4) / rocs_enrich_4)
    except Exception as e:
        print(e)


def plot_hist(data, title="Diff", filename="diff.jpg"):
    import matplotlib.pyplot as plt

    data = [i for i in data if i != np.inf and i != -np.inf]
    fig, ax = plt.subplots(figsize=(8, 6))
    ax.hist(data, color="#807FFF", edgecolor="black", bins=len(data))
    ax.set_xlabel("(ROCS - ROSHAMBO)/ROCS", fontsize=18, fontweight="bold")
    ax.set_ylabel("Frequency", fontsize=18, fontweight="bold")
    # ax.set_yticks(np.arange(0.2, 1.1, 0.2))
    # ax.set_ylim(0, 1.1)
    ax.set_title(title, fontsize=18, fontweight="bold")
    for spine in ["top", "bottom", "left", "right"]:
        ax.spines[spine].set_linewidth(2)
    ax.tick_params(direction="in", labelsize=16, length=6)

    # Save the plot
    plt.savefig(filename, dpi=500, bbox_inches="tight")


# def run_code():
#     sm = GetSimilarityScores(
#         ref_file="query_chembl399440.sdf",
#         dataset_files_pattern="actives_adrb1_458.sdf",
#         opt=False,
#         ignore_hs=True,
#         n_confs=20,
#         keep_mol=True,
#         random_seed=109838974,
#         opt_confs=False,
#         calc_energy=True,
#         energy_iters=300,
#         energy_cutoff=100,
#         align_confs=True,
#         rms_cutoff=0.1,
#         num_threads=28,
#         method="ETKDGv2",
#         working_dir="/UserUCDD/ratwi/tests/20230201_paper_datasets/adrb1/roshambo/actives",
#     )
#     sm.run_paper()
#     sm.convert_transformation_arrays()
#     sm.transform_molecules()
#     sm.calculate_scores(
#         volume_type="analytic",
#         n=2,
#         epsilon=0,
#         use_carbon_radii=True,
#         color=True,
#         max_conformers=1,
#         sort_by="ComboTanimoto",
#         write_to_file=True,
#     )
#
#
# # cProfile.run('run_code()')
# pr = cProfile.Profile()
# pr.enable()
# run_code()
# pr.disable()
# pr.print_stats(sort='tottime')

# st = time.time()
# inputs = [(ref_pharm, pharm) for pharm in outputs_pharm]
# with Pool(processes=cpu_count()) as pool:
#     fit_overlaps = pool.starmap(calc_pharm_overlap, inputs)
# et = time.time()
# print(et - st)
#
# st = time.time()
# for mol in mols:
#     calc_pharmacophore(mol)
# et = time.time()
# print(et - st)


def create_ppt(input_dir, output_file):
    # Create a new PowerPoint presentation
    ppt = Presentation()

    # Define the layout for the slides
    slide_layout = ppt.slide_layouts[5]

    # Iterate through all files in the input directory
    for file in sorted(os.listdir(input_dir)):
        # Check if the file is a non-log plot (exclude log plots)
        if file.endswith(".jpg") and not file.endswith("_log.jpg"):
            # Check if the corresponding log plot exists
            log_file = os.path.splitext(file)[0] + "_log.jpg"
            if os.path.exists(os.path.join(input_dir, log_file)):
                # Add a new slide with the defined layout
                slide = ppt.slides.add_slide(slide_layout)

                # Add the non-log plot
                left = Inches(0.5)
                top = Inches(1.5)
                pic = slide.shapes.add_picture(
                    os.path.join(input_dir, file), left, top, width=Inches(4.5)
                )

                # Add the log plot
                left = Inches(5.5)
                top = Inches(1.5)
                pic = slide.shapes.add_picture(
                    os.path.join(input_dir, log_file), left, top, width=Inches(4.5)
                )

    # Save the PowerPoint presentation
    ppt.save(output_file)


input_dir = "../roshambo"
output_file = "output.pptx"
create_ppt(input_dir, output_file)


def create_folder_list(dir_path):
    folder_list = []
    folder_names = [
        folder_name
        for folder_name in os.listdir(dir_path)
        if os.path.isdir(os.path.join(dir_path, folder_name))
    ]
    for i in range(0, len(folder_names), 7):
        folder_info = tuple(folder_names[i : i + 7])
        folder_list.append(folder_info)
    return folder_list


dir_path = "../roshambo"
folder_list = create_folder_list(dir_path)

print(folder_list)
